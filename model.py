from pptx import Presentation
from pptx.util import Pt
from transformers import AutoTokenizer
from optimum.intel.openvino import OVModelForSeq2SeqLM
import os, re


# ---------------------------
# PPTファイルを読み書きするクラス
# ---------------------------
class PPTModel:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.presentation = Presentation(ppt_path)
        self.edited_ppt_path = None

    def extract_slides_text(self):
        """
        スライドごとに [[shape1_text, shape2_text, ...], [...], ...] の形式で返す
        """
        slides_text = []
        for slide in self.presentation.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    slide_texts.append(shape.text)
            slides_text.append(slide_texts)
        return slides_text

    def update_slide_text(self, slide_idx, new_texts):
        """
        フォーマットを崩さずスライド内のテキストを書き換える
        改行は段落に変換
        """
        if slide_idx < 0 or slide_idx >= len(self.presentation.slides):
            return

        slide = self.presentation.slides[slide_idx]
        shape_idx = 0

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape_idx >= len(new_texts):
                    break

                text_frame = shape.text_frame
                new_text = new_texts[shape_idx].strip()

                if not new_text:
                    shape_idx += 1
                    continue

                # 既存段落を削除（フォーマット崩れを防ぎつつクリア）
                for p in list(text_frame.paragraphs):
                    text_frame._element.remove(p._element)

                # 改行を段落に変換
                for line in new_text.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    p = text_frame.add_paragraph()
                    p.text = line
                    for run in p.runs:
                        run.font.size = Pt(18)  # 必要なら調整

                shape_idx += 1

    def save(self, suffix="_edited"):
        """
        編集後PPTを保存
        """
        self.edited_ppt_path = self.ppt_path.replace(".pptx", f"{suffix}.pptx")
        self.presentation.save(self.edited_ppt_path)
        return self.edited_ppt_path


# ---------------------------
# 翻訳・トーンアップ モデルクラス
# ---------------------------
class TranslatorModel:
    def __init__(self, model_dir="openvino_model", src_lang="ja_XX", tgt_lang="en_XX"):
        print("✅ Loading OpenVINO model…")
        self.tokenizer = AutoTokenizer.from_pretrained(model_dir)
        self.model = OVModelForSeq2SeqLM.from_pretrained(model_dir)

        self.src_lang = src_lang
        self.tgt_lang = tgt_lang
        self.tokenizer.src_lang = src_lang
        self.forced_bos_token_id = self.tokenizer.lang_code_to_id[tgt_lang]
        print("✅ OpenVINO model loaded!")

    def translate_text(self, text: str) -> str:
        if not text.strip():
            return ""

        # 不要な改行を整形
        text = re.sub(r'\n{2,}', '\n', text.strip())

        inputs = self.tokenizer(text, return_tensors="pt", truncation=True, max_length=256)
        outputs = self.model.generate(
            **inputs,
            max_length=256,
            forced_bos_token_id=self.forced_bos_token_id
        )
        return self.tokenizer.decode(outputs[0], skip_special_tokens=True)

    def tone_up(self, text: str, style_prompt: str = "文体を丁寧に、明るく、日本語でかきかえて") -> str:
        if not text.strip():
            return ""

        input_text = f"{style_prompt}: {text}"
        forced_bos_token_id_ja = self.tokenizer.lang_code_to_id["ja_XX"]

        inputs = self.tokenizer(input_text, return_tensors="pt", truncation=True, max_length=256)
        outputs = self.model.generate(
            **inputs,
            max_length=256,
            forced_bos_token_id=forced_bos_token_id_ja
        )
        return self.tokenizer.decode(outputs[0], skip_special_tokens=True)
